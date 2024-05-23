import glob
import io
import json
import re

from docx import Document
from pptx import Presentation


class FileInfoExtractor:
    def __init__(self, file):
        """
        FileInfoExtractor クラス初期化メソッド

        Args:
            file (werkzeug.FileStorage): フォームから受け取ったファイル

        Note:
            ローカルからファイルを読みたいときは、一度 FileStorage で変換する
            from werkzeug.datastructures import FileStorage
            with open('ABC.docx', 'rb') as file:
                file_storage = FileStorage(stream=file, filename='file_name', content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                FileInfoExtractor(file_storage)
        """
        self.file = file
        if file.mimetype == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            self.type = 'docx'
        elif file.mimetype == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            self.type = 'pptx'

        self._file_data_array = self.__extract()

    @property
    def file_data_array(self):
        """
        docxやpptxなどのファイルのデータをlistにする
        """
        return self._file_data_array

    @property
    def file_data_dict(self):
        """
        docxやpptxなどのファイルのデータをdictにする
        """
        return {
            'name': self.file.filename,
            'type': self.type,
            'contents': self._file_data_array
        }

    def __extract(self):
        """
        ファイルからデータを抽出する

        Returns:
            list: データのlist
        """
        if self.type == 'docx':
            return self.__extract_docx()
        elif self.type == 'pptx':
            return self.__extract_pptx()
    
    def __extract_docx(self):
        """
        docxファイルからデータを抽出する

        Returns:
            list: データのlist
        """
        file_stream = io.BytesIO(self.file.read())
        doc = Document(file_stream)

        texts = ""
        for para in doc.paragraphs:
            text = re.sub(r"\s", "", para.text)
            if text != "":
                texts += text
            else:
                texts += "\n"

        return texts.replace("株式会社イーネット","")

    def __extract_pptx(self):
        """
        Extracts data from a pptx file excluding content from tables.

        Returns:
            list: データのlist
        """
        file_stream = io.BytesIO(self.file.read())
        pptx = Presentation(file_stream)
        texts = []

        for slide in pptx.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and not shape.has_table:
                    text = shape.text.strip().replace("\n", " ")
                    if len(text) > 15:
                        texts.append(text)

        pptx_content = "\n".join(texts)

        return pptx_content