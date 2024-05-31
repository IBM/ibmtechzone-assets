from pptx import Presentation
import io
from docx import Document
import io
from werkzeug.datastructures import FileStorage
from pptx import Presentation

class FileInfoExtractor:
    def __init__(self, file):
        """
        FileInfoExtractor 

        Args:
            file (werkzeug.FileStorage):

        Note:
            If you want to read a file locally, convert it using FileStorage once. 
            
            from werkzeug.datastructures import FileStorage
            with open('ABC.pptx', 'rb') as file:
                file_storage = FileStorage(stream=file, filename='file_name', content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
                FileInfoExtractor(file_storage)
        """
        self.file = file
        if file.mimetype == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            self.type = 'docx'
        elif file.mimetype == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            self.type = 'pptx'

        self._file_data_array = self.__extract()

    @property
    def file_data_array(self):
        """
        List data from files such as docx and pptx        
        """
        return self._file_data_array

    @property
    def file_data_dict(self):
        """
        Convert file data such as docx and pptx to dict
        """
        return {
            'name': self.file.filename,
            'type': self.type,
            'contents': self._file_data_array
        }

    def __extract(self):
        """
        Extract data from file

        Returns:
            list: list of data
        """
        if self.type == 'docx':
            return self.__extract_docx()
        elif self.type == 'pptx':
            return self.__extract_pptx()
    
    def __extract_docx(self):
        """
        Extract data from docx file

        Returns:
            list: list of data
        """
        file_stream = io.BytesIO(self.file.read())
        doc = Document(file_stream)

        texts = ""
        for para in doc.paragraphs:
            text = re.sub(r"\s", "", para.text)
            if text != "":
                texts += text
            else:
                texts += "\n"

        return texts

    def __extract_pptx(self):
        """
        Extracts data from a pptx file including content from tables, preserving the row-column relationship.

        Returns:
            str: concatenated text data
        """
        file_stream = io.BytesIO(self.file.read())
        pptx = Presentation(file_stream)
        texts = []

        # for slide in pptx.slides:
        #     for shape in slide.shapes:
        #         if hasattr(shape, "text") and not shape.has_table:
        #             text = shape.text.strip().replace("\n", " ")
        #             if len(text) > 15:
        #                 texts.append(text)

        # pptx_content = "\n".join(texts)

        # return pptx_content

        for slide in pptx.slides:
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_table:
                    table_data = []
                    tbl = shape.table
                    for row in tbl.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = []
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    cell_text.append(run.text)
                            row_data.append(" ".join(cell_text))
                        table_data.append(row_data)
                    slide_texts.append(table_data)
                elif hasattr(shape, "text_frame") and shape.text_frame is not None:
                    text = shape.text_frame.text.strip().replace("\n", " ")
                    if len(text) > 15:
                        slide_texts.append(text)
            texts.append(slide_texts)

        return texts    

def main():
    path = "/Users/kanishksaxena/Documents/POC's/generate-approval-docs/input documents/proposal_document.pptx"
        # Open the file in binary mode and wrap it in a FileStorage object
    with open(path, 'rb') as pptx_file:
        pptx_file_storage = FileStorage(
            stream=pptx_file,
            filename='proposal_document.pptx',
            content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
        # Create an instance of FileInfoExtractor with the PPTX file
        file_extractor = FileInfoExtractor(pptx_file_storage)
        
        # Extract and print the text data from the PPTX file
        pptx_text = file_extractor.file_data_array
        print("Extracted PPTX Text:")
        print(pptx_text)

if __name__ == "__main__":
    main()
